"""知识库构建模块：四级自适应分块 + 父子分块（双层检索）

分块策略参考 WeKnora 文档分块指南：
- auto（默认）：自动分析文档特征，择优匹配策略
- heading：按 Markdown 标题边界切割，拼接面包屑上下文
- heuristic：PDF 结构化文档，按分页符/编号章节/视觉分隔线切割
- legacy：纯递归分块（兜底）

核心参数：512 字符 + 80 字符重叠（Vecta 基准推荐）
计量方式：字符计数（与 WeKnora/Vecta 一致）
"""

import os
import re
from typing import List, Optional

import fitz  # pymupdf

from backend.core.config import settings


_knowledge_base: "KnowledgeBase | None" = None


def get_knowledge_base() -> "KnowledgeBase":
    """获取知识库单例"""
    global _knowledge_base
    if _knowledge_base is None:
        _knowledge_base = KnowledgeBase()
    return _knowledge_base


class KnowledgeBase:
    """知识库管理 — 四级自适应分块 + 父子双层检索"""

    def __init__(self):
        self.chunk_size = settings.RAG_CHUNK_SIZE
        self.chunk_overlap = settings.RAG_CHUNK_OVERLAP
        self.strategy = settings.RAG_CHUNK_STRATEGY
        self.separators = settings.RAG_CHUNK_SEPARATORS
        self.embedding_token_limit = settings.RAG_EMBEDDING_TOKEN_LIMIT
        self.parent_child_enabled = settings.RAG_PARENT_CHILD_ENABLED
        self.parent_chunk_size = settings.RAG_PARENT_CHUNK_SIZE
        self.child_chunk_size = settings.RAG_CHILD_CHUNK_SIZE
        self.knowledge_dir = settings.RAG_KNOWLEDGE_DIR
        os.makedirs(self.knowledge_dir, exist_ok=True)

    # ============================================================
    #  文档解析
    # ============================================================

    def parse_pdf(self, file_path: str) -> str:
        sections = self.parse_pdf_pages(file_path)
        return "\n".join(s['text'] for s in sections)

    def parse_pdf_pages(self, file_path: str) -> List[dict]:
        """解析PDF，按页提取文本"""
        sections = []
        try:
            doc = fitz.open(file_path)
            for i, page in enumerate(doc, 1):
                page_text = page.get_text()
                if page_text:
                    cleaned = self._clean_text(page_text)
                    if cleaned:
                        sections.append({'text': cleaned, 'page': i})
            doc.close()
        except Exception as e:
            print(f"PDF解析失败: {e}")
        return sections

    def parse_txt(self, file_path: str) -> str:
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return self._clean_text(f.read())
        except Exception as e:
            print(f"TXT解析失败: {e}")
            return ""

    def parse_md(self, file_path: str) -> str:
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()
            return re.sub(r'\n{3,}', '\n\n', text)
        except Exception as e:
            print(f"Markdown解析失败: {e}")
            return ""

    def parse_docx(self, file_path: str) -> str:
        try:
            from docx import Document
            doc = Document(file_path)
            parts = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
            for table in doc.tables:
                for row in table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        parts.append(" | ".join(cells))
            return self._clean_text("\n\n".join(parts))
        except Exception as e:
            print(f"DOCX解析失败: {e}")
            return ""

    def parse_pptx(self, file_path: str) -> str:
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            slides_text = []
            for i, slide in enumerate(prs.slides, 1):
                texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t:
                                texts.append(t)
                    if shape.has_table:
                        for row in shape.table.rows:
                            cells = [c.text.strip() for c in row.cells if c.text.strip()]
                            if cells:
                                texts.append(" | ".join(cells))
                if texts:
                    slides_text.append(f"[第{i}页]\n" + "\n".join(texts))
            return self._clean_text("\n\n".join(slides_text))
        except Exception as e:
            print(f"PPTX解析失败: {e}")
            return ""

    def _clean_text(self, text: str) -> str:
        text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]', '', text)
        text = re.sub(r'[ \t]{2,}', ' ', text)
        text = re.sub(r'\n{3,}', '\n\n', text)
        return text.strip()

    # ============================================================
    #  四级自适应分块策略
    # ============================================================

    def _analyze_document(self, text: str) -> dict:
        """分析文档结构特征，返回信号量"""
        signals = {
            'heading_count': len(re.findall(r'^#{1,6}\s', text, re.MULTILINE)),
            'page_break_count': text.count('\x0c'),
            'numbered_section_count': len(re.findall(r'^\d+[\.\)]\s', text, re.MULTILINE)),
            'all_caps_title_count': len(re.findall(r'^[A-Z][A-Z\s]{5,}$', text, re.MULTILINE)),
            'separator_line_count': len(re.findall(r'^[-=_*]{5,}$', text, re.MULTILINE)),
            'empty_line_ratio': text.count('\n\n') / max(len(text), 1),
            'total_chars': len(text),
        }
        return signals

    def _select_strategy(self, text: str, override: str = None) -> str:
        """根据文档信号选择最优分块策略"""
        if override and override != "auto":
            return override

        signals = self._analyze_document(text)

        # 标题分块：Markdown 标题 ≥ 3 个，且标题数合理（非每个短行都是标题）
        if signals['heading_count'] >= 3 and signals['heading_count'] < len(text.split('\n')) * 0.3:
            return "heading"

        # 启发式分块：有分页符或编号章节
        if signals['page_break_count'] >= 2 or signals['numbered_section_count'] >= 3:
            return "heuristic"

        # 兜底：递归分块
        return "legacy"

    def split_into_chunks(self, text: str, strategy: str = None) -> List[str]:
        """将文本分块（无页码信息）"""
        return [c['content'] for c in self.split_into_chunks_with_pages(
            [{'text': text, 'page': None}], strategy=strategy
        )]

    def split_into_chunks_with_pages(
        self, sections: List[dict], strategy: str = None
    ) -> List[dict]:
        """将带页码的文本分块

        Args:
            sections: [{'text': str, 'page': int|None}, ...]
            strategy: auto/heading/heuristic/legacy，None 则用配置默认值

        Returns:
            [{'content': str, 'page': int|None, 'strategy': str}, ...]
        """
        use_strategy = strategy or self.strategy
        all_text = "\n".join(s.get('text', '') for s in sections)
        selected = self._select_strategy(all_text, override=use_strategy)

        if selected == "heading":
            chunks = self._split_by_heading(sections)
        elif selected == "heuristic":
            chunks = self._split_by_heuristic(sections)
        else:
            chunks = self._split_by_legacy(sections)

        # 降级校验：如果策略生成异常多的极短分块，降级到 legacy
        if selected != "legacy" and len(chunks) > 0:
            short_chunks = sum(1 for c in chunks if len(c['content']) < 30)
            if short_chunks > len(chunks) * 0.5:
                print(f"策略 {selected} 生成过多短分块 ({short_chunks}/{len(chunks)})，降级到 legacy")
                chunks = self._split_by_legacy(sections)
                selected = "legacy"

        # 添加重叠
        if self.chunk_overlap > 0 and len(chunks) > 1:
            overlapped = []
            for i, chunk in enumerate(chunks):
                content = chunk['content']
                if i > 0:
                    prev = chunks[i - 1]['content']
                    overlap_text = prev[-self.chunk_overlap:] if len(prev) > self.chunk_overlap else prev
                    content = overlap_text + content
                overlapped.append({
                    'content': content,
                    'page': chunk.get('page'),
                    'strategy': selected,
                })
            chunks = overlapped
        else:
            for c in chunks:
                c['strategy'] = selected

        return chunks

    # --- heading 标题分块 ---
    def _split_by_heading(self, sections: List[dict]) -> List[dict]:
        """按 Markdown 标题边界切割，拼接面包屑上下文"""
        chunks = []
        current_content = ""
        current_page = None
        breadcrumb = []  # 面包屑：[("##", "章节名"), ...]

        for section in sections:
            text = section.get('text', '')
            page = section.get('page')
            if not text:
                continue

            lines = text.split('\n')
            for line in lines:
                heading_match = re.match(r'^(#{1,6})\s+(.+)$', line)
                if heading_match:
                    # 遇到标题：先保存当前块
                    if current_content.strip():
                        prefix = " > ".join(f"{h[1]}" for h in breadcrumb) + "\n" if breadcrumb else ""
                        chunks.append({
                            'content': prefix + current_content.strip(),
                            'page': current_page,
                        })

                    # 更新面包屑
                    level = len(heading_match.group(1))
                    title = heading_match.group(2).strip()
                    breadcrumb = [(l, t) for l, t in breadcrumb if l < level]
                    breadcrumb.append((level, title))

                    current_content = line + "\n"
                    current_page = page
                else:
                    current_content += line + "\n"
                    if current_page is None and page is not None:
                        current_page = page

                    # 超长检查
                    if len(current_content) > self.chunk_size:
                        prefix = " > ".join(f"{h[1]}" for h in breadcrumb) + "\n" if breadcrumb else ""
                        chunks.append({
                            'content': prefix + current_content.strip(),
                            'page': current_page,
                        })
                        current_content = ""
                        current_page = None

        if current_content.strip():
            prefix = " > ".join(f"{h[1]}" for h in breadcrumb) + "\n" if breadcrumb else ""
            chunks.append({
                'content': prefix + current_content.strip(),
                'page': current_page,
            })

        return chunks

    # --- heuristic 启发式分块 ---
    def _split_by_heuristic(self, sections: List[dict]) -> List[dict]:
        """依据分页符、编号章节、全大写标题、分隔线等特征智能切割"""
        chunks = []
        current_content = ""
        current_page = None

        for section in sections:
            text = section.get('text', '')
            page = section.get('page')
            if not text:
                continue

            # 按启发式特征分割
            parts = re.split(
                r'(\n\x0c\n|'           # 分页符
                r'\n\d+[\.\)]\s.+?\n|'   # 编号章节 (1. xxx / 1) xxx)
                r'\n[A-Z][A-Z\s]{5,}\n|' # 全大写标题
                r'\n[-=_*]{5,}\n)',       # 视觉分隔线
                text
            )

            for part in parts:
                if not part or not part.strip():
                    continue
                # 跳过分隔符本身
                if re.match(r'^\s*[\x0c\-=_*]+$', part):
                    # 分隔符处断开
                    if current_content.strip():
                        chunks.append({'content': current_content.strip(), 'page': current_page})
                        current_content = ""
                        current_page = None
                    continue

                if len(current_content) + len(part) <= self.chunk_size:
                    current_content += part
                    if current_page is None and page is not None:
                        current_page = page
                else:
                    if current_content.strip():
                        chunks.append({'content': current_content.strip(), 'page': current_page})
                    current_content = part
                    current_page = page

        if current_content.strip():
            chunks.append({'content': current_content.strip(), 'page': current_page})

        return chunks

    # --- legacy 递归分块 ---
    def _split_by_legacy(self, sections: List[dict]) -> List[dict]:
        """纯递归分块：按分隔符优先级逐级切割"""
        chunks = []
        current_content = ""
        current_page = None

        for section in sections:
            text = section.get('text', '')
            page = section.get('page')
            if not text:
                continue

            # 先按段落分割
            paragraphs = text.split('\n\n')
            for para in paragraphs:
                para = para.strip()
                if not para:
                    continue

                if len(current_content) + len(para) + 2 <= self.chunk_size:
                    current_content += "\n\n" + para if current_content else para
                    if current_page is None and page is not None:
                        current_page = page
                else:
                    if current_content:
                        chunks.append({'content': current_content, 'page': current_page})
                        current_content = ""
                        current_page = None

                    if len(para) > self.chunk_size:
                        # 按分隔符优先级递归切割
                        sub_chunks = self._recursive_split(para, self.separators)
                        for sc in sub_chunks:
                            if len(current_content) + len(sc) + 2 <= self.chunk_size:
                                current_content += "\n\n" + sc if current_content else sc
                                if current_page is None and page is not None:
                                    current_page = page
                            else:
                                if current_content:
                                    chunks.append({'content': current_content, 'page': current_page})
                                current_content = sc
                                current_page = page
                    else:
                        current_content = para
                        if current_page is None and page is not None:
                            current_page = page

        if current_content:
            chunks.append({'content': current_content, 'page': current_page})

        return chunks

    def _recursive_split(self, text: str, separators: list) -> List[str]:
        """按分隔符优先级递归切割"""
        if len(text) <= self.chunk_size:
            return [text]
        if not separators:
            # 硬切
            return [text[i:i + self.chunk_size] for i in range(0, len(text), self.chunk_size)]

        sep = separators[0]
        remaining_seps = separators[1:]

        parts = text.split(sep)
        result = []
        current = ""
        for part in parts:
            if not part:
                continue
            if len(current) + len(part) + len(sep) <= self.chunk_size:
                current += sep + part if current else part
            else:
                if current:
                    result.append(current)
                if len(part) > self.chunk_size:
                    result.extend(self._recursive_split(part, remaining_seps))
                    current = ""
                else:
                    current = part
        if current:
            result.append(current)
        return result

    # ============================================================
    #  父子分块（双层检索）
    # ============================================================

    def split_into_parent_child(self, text: str, strategy: str = None) -> dict:
        """父子分块：返回父分块和子分块

        Returns:
            {
                'parents': [{'content': str, 'index': int}, ...],
                'children': [{'content': str, 'parent_index': int, 'page': int|None}, ...],
            }
        """
        if not self.parent_child_enabled:
            # 未启用父子分块，退化为单层
            chunks = self.split_into_chunks(text, strategy=strategy)
            pseudo_parents = [{'content': c, 'index': i} for i, c in enumerate(chunks)]
            pseudo_children = [{'content': c, 'parent_index': i, 'page': None} for i, c in enumerate(chunks)]
            return {'parents': pseudo_parents, 'children': pseudo_children}

        # 先做父分块（大尺寸）
        old_size = self.chunk_size
        self.chunk_size = self.parent_chunk_size
        parent_chunks = self.split_into_chunks(text, strategy=strategy)
        self.chunk_size = old_size

        parents = []
        children = []

        for p_idx, p_chunk in enumerate(parent_chunks):
            parents.append({'content': p_chunk, 'index': p_idx})

            # 将每个父分块切为子分块
            child_size = self.child_chunk_size
            child_overlap = int(child_size * 0.15)
            text_p = p_chunk
            p_len = len(text_p)

            if p_len <= child_size:
                children.append({
                    'content': text_p,
                    'parent_index': p_idx,
                    'page': None,
                })
            else:
                i = 0
                while i < p_len:
                    end = min(i + child_size, p_len)
                    child_text = text_p[i:end]
                    children.append({
                        'content': child_text,
                        'parent_index': p_idx,
                        'page': None,
                    })
                    i += child_size - child_overlap
                    if i >= p_len:
                        break

        return {'parents': parents, 'children': children}

    # ============================================================
    #  文件处理入口
    # ============================================================

    def process_file(self, file_path: str) -> List[str]:
        return [c['content'] for c in self.process_file_with_metadata(file_path)]

    def process_file_with_metadata(self, file_path: str) -> List[dict]:
        """处理文件，返回带元数据的文本块

        Returns: [{'content': str, 'page': int|None, 'strategy': str}, ...]
        """
        ext = os.path.splitext(file_path)[1].lower()

        if ext == '.pdf':
            sections = self.parse_pdf_pages(file_path)
            if not sections:
                return []
            return self.split_into_chunks_with_pages(sections)
        elif ext == '.pptx':
            text = self.parse_pptx(file_path)
            if not text:
                return []
            sections = self._pptx_to_sections(text)
            return self.split_into_chunks_with_pages(sections)
        elif ext in ('.txt', '.md', '.docx'):
            if ext == '.txt':
                text = self.parse_txt(file_path)
            elif ext == '.md':
                text = self.parse_md(file_path)
            else:
                text = self.parse_docx(file_path)
            if not text:
                return []
            return self.split_into_chunks_with_pages([{'text': text, 'page': None}])
        else:
            print(f"不支持的文件类型: {ext}")
            return []

    def _pptx_to_sections(self, text: str) -> List[dict]:
        sections = []
        parts = re.split(r'\[第(\d+)页\]\n?', text)
        i = 1
        while i < len(parts):
            page_num = int(parts[i])
            slide_text = parts[i + 1] if i + 1 < len(parts) else ''
            if slide_text.strip():
                sections.append({'text': slide_text.strip(), 'page': page_num})
            i += 2
        return sections

    def index_history_reports(self, reports: List[str]) -> List[str]:
        all_chunks = []
        for report in reports:
            if report:
                chunks = self.split_into_chunks(report)
                all_chunks.extend(chunks)
        return all_chunks

    # ============================================================
    #  分块预览
    # ============================================================

    def preview_chunks(self, text: str, strategy: str = None) -> dict:
        """预览分块结果（不写入数据库/索引）

        Returns:
            {
                'strategy': str,           # 实际使用的策略
                'strategy_source': str,    # auto 选择的依据
                'total_chunks': int,
                'stats': {
                    'avg_chars': float,
                    'min_chars': int,
                    'max_chars': int,
                    'std_chars': float,
                },
                'signals': dict,           # 文档结构信号
                'chunks': [{'index': int, 'content': str, 'chars': int, 'page': ...}, ...],
            }
        """
        signals = self._analyze_document(text)
        use_strategy = strategy or self.strategy
        selected = self._select_strategy(text, override=use_strategy)

        chunks = self.split_into_chunks_with_pages(
            [{'text': text, 'page': None}], strategy=use_strategy
        )

        chunk_sizes = [len(c['content']) for c in chunks]
        import statistics
        stats = {
            'avg_chars': round(statistics.mean(chunk_sizes), 1) if chunk_sizes else 0,
            'min_chars': min(chunk_sizes) if chunk_sizes else 0,
            'max_chars': max(chunk_sizes) if chunk_sizes else 0,
            'std_chars': round(statistics.stdev(chunk_sizes), 1) if len(chunk_sizes) > 1 else 0,
        }

        return {
            'strategy': selected,
            'strategy_source': f"auto → {selected}" if use_strategy == "auto" else f"manual → {selected}",
            'total_chunks': len(chunks),
            'stats': stats,
            'signals': signals,
            'chunks': [
                {
                    'index': i,
                    'content': c['content'][:200] + ('...' if len(c['content']) > 200 else ''),
                    'chars': len(c['content']),
                    'page': c.get('page'),
                }
                for i, c in enumerate(chunks)
            ],
        }
